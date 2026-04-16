from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = r"D:\work\15_지선차단기\09_보고서\protocol_agenda_approval_2026-04-14.pptx"

TITLE_COLOR = RGBColor(31, 51, 76)
ACCENT = RGBColor(29, 78, 137)
LIGHT = RGBColor(241, 246, 252)
LIGHT2 = RGBColor(250, 252, 255)
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)


def style_run(run, size=Pt(12), bold=False, color=TEXT):
    run.font.name = "맑은 고딕"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle="SEMS / 추가 논의 안건 / 2026-04-14"):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_COLOR
    band.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(11.8), Inches(0.45))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, Pt(24), True, TITLE_COLOR)

    sbox = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(5.5), Inches(0.25))
    p = sbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    style_run(r, Pt(10), False, MUTED)


def add_box(slide, left, top, width, height, title, lines, fill_color=WHITE, title_color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(210, 220, 232)
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, Pt(16), True, title_color)

    for idx, line in enumerate(lines):
        p = tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4 if idx == 0 else 1)
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = f"- {line}"
        style_run(r, Pt(11), False, TEXT)


def make_slide(slide, agenda_title, conclusion, reasons):
    add_title(slide, agenda_title)
    add_box(
        slide,
        Inches(0.65), Inches(1.75), Inches(12.0), Inches(1.2),
        "결론",
        [conclusion],
        fill_color=LIGHT,
        title_color=TITLE_COLOR,
    )
    add_box(
        slide,
        Inches(0.65), Inches(3.15), Inches(12.0), Inches(3.35),
        "사유",
        reasons,
        fill_color=LIGHT2,
        title_color=ACCENT,
    )


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
make_slide(
    slide,
    "안건 1. 2차 개발에 이벤트 기능을 추가할지 여부",
    "2차 개발에 이벤트 기능을 추가하는 방향으로 진행",
    [
        "과전압, 과전류는 보호 기능 성격이므로 가능한 한 빠르게 릴레이가 차단되어야 함",
        "DCU가 polling으로 전압, 전류를 읽은 뒤 판단해서 다시 SL에 차단 명령을 보내는 방식은 반응이 늦어질 수 있음",
        "따라서 SL이 현장에서 즉시 판단하고 릴레이를 OFF한 뒤, 그 결과를 이벤트로 DCU에 전달하는 구조가 더 적절함",
        "이 방식이 보호 속도와 이벤트 이력 관리 측면에서 유리함",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
make_slide(
    slide,
    "안건 2. 이벤트 기능을 넣는다면 DCU를 어떤 방식으로 개발할지 여부",
    "추가선 없이 기존 RS-485 통신선에서 비동기 이벤트 방식으로 개발",
    [
        "별도 배선을 추가하지 않고도 UART/RS-485 통신선으로 이벤트 송신이 가능함",
        "평상시에는 DCU polling 방식으로 운용하고, 이상 발생 시에는 SL이 비동기 이벤트를 송신하는 혼합 구조로 개발 가능함",
        "이 방식은 배선 증가 없이도 이상 상황을 빠르게 전달할 수 있고, 서버 보고 및 장애 이력 관리에도 적합함",
        "추가 하드웨어 변경을 최소화하면서 2차 기능 확장이 가능함",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
make_slide(
    slide,
    "안건 3. 통신 속도를 몇으로 적용할지 여부",
    "초기 적용 속도는 9600 bps로 진행",
    [
        "현재 예상되는 SL 연결 대수가 10대 내외 수준이므로 고속 통신보다는 안정성을 우선하는 것이 적절함",
        "RS-485 환경에서는 속도를 무리하게 높이기보다 노이즈 내성과 통신 신뢰성을 확보하는 것이 중요함",
        "초기 버전은 9600 bps로 안정적으로 운용하고, 향후 연결 대수 증가나 데이터량 증가 시 속도 상향을 재검토하는 것이 합리적임",
    ],
)

prs.save(OUT_PATH)

with ZipFile(OUT_PATH) as zf:
    xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "안건 1. 2차 개발에 이벤트 기능을 추가할지 여부" in xml
    assert "결론" in xml
    assert "사유" in xml

print(OUT_PATH)
print("VERIFY_OK")
