from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = r"D:\work\15_지선차단기\09_보고서\protocol_full_spec_tables_v4_2026-04-14.pptx"

TITLE_COLOR = RGBColor(31, 51, 76)
ACCENT = RGBColor(29, 78, 137)
HEADER_FILL = RGBColor(225, 235, 247)
SUB_FILL = RGBColor(240, 245, 251)
BORDER = RGBColor(180, 190, 205)
TEXT = RGBColor(40, 40, 40)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=Pt(12), bold=False, color=TEXT):
    run.font.name = "맑은 고딕"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle="SEMS / Protocol / 2026-04-14"):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_COLOR
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(7.5), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, Pt(24), True, TITLE_COLOR)

    sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.28), Inches(4.5), Inches(0.3))
    p = sub_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    set_font(run, Pt(10), False, RGBColor(110, 110, 110))


def add_table(slide, left, top, width, height, headers, rows, font_size=Pt(11), first_col_bold=False):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    col_width = int(width / len(headers))
    for idx in range(len(headers)):
        table.columns[idx].width = col_width

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_font(run, font_size, True, TITLE_COLOR)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else SUB_FILL
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    set_font(run, font_size, first_col_bold and c == 0)

    for row in table.rows:
        row.height = int(height / (len(rows) + 1))

    return table


def add_note(slide, text, left=0.55, top=6.7, width=12.1, height=0.45):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, Pt(11), False, TEXT)


def add_bullets(slide, title, bullets, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, Pt(14), True, ACCENT)
    for bullet in bullets:
        p = tf.add_paragraph()
        p.level = 0
        p.text = f"- {bullet}"
        for run in p.runs:
            set_font(run, Pt(11))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "지선차단기 - DCU 통신 프로토콜", "프레임 형식 / 필드 정의")
add_table(
    slide,
    Inches(0.45),
    Inches(1.7),
    Inches(12.4),
    Inches(3.3),
    ["필드", "크기", "설명", "초기값 / 비고"],
    [
        ["STX", "1 byte", "프레임 시작 표시", "0x7E"],
        ["VER", "1 byte", "프로토콜 버전", "0x01"],
        ["MSG_TYPE", "1 byte", "메시지 종류", "Read / Control / ACK / NACK"],
        ["SEQ", "1 byte", "요청-응답 식별 번호", "0x00 ~ 0xFF 순환"],
        ["SL_ID", "1 byte", "장치 식별자", "장치별 고유 ID"],
        ["LEN", "2 byte", "Payload 길이", "byte 단위"],
        ["PAYLOAD", "N byte", "실제 데이터", "TLV 유사 구조 사용"],
        ["CRC16", "2 byte", "프레임 오류 검출", "무결성 확인"],
        ["ETX", "1 byte", "프레임 종료 표시", "0x7E"],
    ],
    font_size=Pt(10.5),
    first_col_bold=True,
)
add_bullets(
    slide,
    "필드 해석 요점",
    [
        "SEQ는 요청-응답 매칭과 재전송 중복 방지를 위해 사용",
        "CRC16은 수신 데이터 손상 검출용",
        "STX/ETX는 프레임 경계 식별용",
    ],
    Inches(0.55),
    Inches(5.25),
    Inches(6.1),
    Inches(1.05),
)
add_note(slide, "프레임 시작 / 종료 바이트는 0x7E로 통일")


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "메시지 종류 및 Payload 구조", "메시지 코드 / TLV 구조 / 3개 단위 전송")
add_table(
    slide,
    Inches(0.5),
    Inches(1.7),
    Inches(5.8),
    Inches(2.55),
    ["코드", "이름", "설명"],
    [
        ["0x01", "Read Request", "데이터 조회 요청"],
        ["0x02", "Read Response", "데이터 조회 응답"],
        ["0x03", "Control Request", "릴레이 제어 요청"],
        ["0x04", "Control Response", "릴레이 제어 응답"],
        ["0x05", "Event Report", "2차 확장용 이벤트 보고"],
        ["0x06", "ACK", "정상 수신 확인"],
        ["0x07", "NACK", "오류 또는 거부 응답"],
    ],
    font_size=Pt(10.5),
)
add_table(
    slide,
    Inches(6.55),
    Inches(1.7),
    Inches(6.15),
    Inches(2.15),
    ["필드", "크기", "설명"],
    [
        ["ITEM_ID", "2 byte", "데이터 항목 ID"],
        ["TYPE", "1 byte", "데이터 타입"],
        ["ITEM_LEN", "1 byte", "VALUE 길이"],
        ["VALUE", "N byte", "실제 값"],
    ],
    font_size=Pt(10.5),
    first_col_bold=True,
)
add_bullets(
    slide,
    "전송 규칙",
    [
        "응답 Payload는 최대 3개 항목 단위로 전송",
        "4개 이상이면 3개 + 나머지 형태로 분할 전송",
        "남는 슬롯은 data_empty(0x0000)로 패딩",
    ],
    Inches(6.6),
    Inches(4.1),
    Inches(5.9),
    Inches(1.2),
)
add_note(slide, "예: 4개 항목 응답 시 1프레임 3개 + 2프레임 1개 + data_empty + data_empty")


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "데이터 항목 정의", "지선차단기 -> DCU / DCU -> 지선차단기")
add_table(
    slide,
    Inches(0.35),
    Inches(1.6),
    Inches(12.6),
    Inches(2.55),
    ["항목 ID", "항목명", "타입", "단위", "필수 여부", "설명"],
    [
        ["0x0000", "data_empty", "uint32", "-", "패딩", "빈 슬롯 채움"],
        ["0x0001", "SL_ID", "uint8", "-", "필수", "장치 식별자"],
        ["0x0002", "Voltage", "uint32", "0.1 V", "필수", "실효 전압"],
        ["0x0003", "Current", "uint32", "0.01 A", "필수", "실효 전류"],
        ["0x0004", "Energy_Acc", "uint32", "Wh", "필수", "누적 전력량"],
        ["0x0005", "Relay_State", "uint8", "-", "필수", "릴레이 현재 상태"],
        ["0x0006", "Alarm_Code", "uint16 bit field", "-", "권장", "이상 상태 코드"],
    ],
    font_size=Pt(10),
)
add_table(
    slide,
    Inches(0.85),
    Inches(4.45),
    Inches(11.6),
    Inches(1.55),
    ["항목 ID", "항목명", "타입", "필수 여부", "설명"],
    [
        ["0x1001", "Read_Item_List", "list", "필수", "조회할 항목 목록"],
        ["0x1002", "Relay_Command", "uint8", "필수", "릴레이 ON/OFF 명령"],
        ["0x1003", "Control_Reason", "uint8", "권장", "제어 사유 코드"],
        ["0x1004", "Control_Token", "uint32", "권장", "제어 명령 검증값"],
    ],
    font_size=Pt(10.5),
)
add_note(slide, "전압 / 전류 / 전력량은 4 byte로 통일")


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "타입, 상태값, 알람 정의", "TYPE / Relay State / Alarm Code")
add_table(
    slide,
    Inches(0.45),
    Inches(1.7),
    Inches(3.25),
    Inches(2.3),
    ["TYPE", "의미"],
    [
        ["0x01", "uint8"],
        ["0x02", "uint16"],
        ["0x03", "uint32"],
        ["0x04", "int32"],
        ["0x05", "bit field"],
        ["0x06", "list"],
    ],
    font_size=Pt(10.5),
)
add_table(
    slide,
    Inches(4.0),
    Inches(1.7),
    Inches(3.55),
    Inches(1.95),
    ["값", "의미"],
    [
        ["0x00", "OFF / 차단"],
        ["0x01", "ON / 투입"],
        ["0x02", "전환 중"],
        ["0x03", "이상 또는 판별 불가"],
    ],
    font_size=Pt(10.5),
)
add_table(
    slide,
    Inches(7.85),
    Inches(1.7),
    Inches(4.95),
    Inches(3.15),
    ["비트", "의미"],
    [
        ["bit0", "과전류"],
        ["bit1", "과전압"],
        ["bit2", "저전압"],
        ["bit3", "릴레이 구동 실패"],
        ["bit4", "측정칩 이상"],
        ["bit5", "내부 메모리 이상"],
        ["bit6", "센서 또는 측정값 이상"],
        ["bit7~15", "예약 / 향후 확장"],
    ],
    font_size=Pt(10),
)
add_bullets(
    slide,
    "판정 기준",
    [
        "과전압: 220V 기준 ±10%",
        "과전류: N회 측정값의 Max보다 M% 높은 경우",
        "N, M 값은 별도 확정 필요",
    ],
    Inches(0.6),
    Inches(4.45),
    Inches(5.8),
    Inches(1.05),
)
add_note(slide, "Alarm_Code는 16비트 bit field 방식 사용")


# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "응답 코드 및 통신 정책", "Response Code / Retry / Event / 속도")
add_table(
    slide,
    Inches(0.45),
    Inches(1.7),
    Inches(6.1),
    Inches(2.85),
    ["코드", "의미"],
    [
        ["0x00", "성공"],
        ["0x01", "알 수 없는 메시지"],
        ["0x02", "지원하지 않는 항목"],
        ["0x03", "길이 오류"],
        ["0x04", "CRC 오류"],
        ["0x05", "잘못된 값"],
        ["0x06", "제어 거부"],
        ["0x07", "내부 처리 실패"],
        ["0x08", "토큰 또는 인증 오류"],
    ],
    font_size=Pt(10.5),
)
add_bullets(
    slide,
    "통신 정책",
    [
        "응답 없음 시 최대 3회 재전송",
        "3회 응답 수신 실패 시 서버 Event 발생",
        "SEQ는 1 byte 순환값(0x00 ~ 0xFF)",
        "동일 SEQ 제어 요청 재수신 시 중복 요청으로 처리",
        "통신 속도: 9600 bps",
    ],
    Inches(6.9),
    Inches(1.8),
    Inches(5.35),
    Inches(2.2),
)
add_bullets(
    slide,
    "2차 개발 방향",
    [
        "Event Report(0x05) 기능 추가 검토",
        "추가선 없이 RS-485 비동기 이벤트 방식 검토",
        "평상시 Polling + 이상 시 Event 혼합 구조",
    ],
    Inches(6.9),
    Inches(4.2),
    Inches(5.35),
    Inches(1.55),
)
add_note(slide, "보호 동작은 SL 현장 판단, DCU는 수집 / 전송 / 이력 관리 역할")


# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "예시 프레임", "조회 / 제어 예시")
add_bullets(
    slide,
    "조회 요청 예시",
    [
        "[7E][01][01][21][11][0006][0002 0003 0005][CRC16][7E]",
        "VER=01 / MSG_TYPE=Read Request / SEQ=21 / SL_ID=11",
        "전압, 전류, 릴레이 상태 조회",
    ],
    Inches(0.55),
    Inches(1.75),
    Inches(5.9),
    Inches(1.5),
)
add_bullets(
    slide,
    "조회 응답 예시",
    [
        "[7E][01][02][21][11][0012][0002 03 04 00000898][0003 03 04 00001388][0005 01 01 01][CRC16][7E]",
        "Voltage=220.0V / Current=50.00A / Relay_State=ON",
        "각 항목은 ITEM_ID + TYPE + LEN + VALUE 구조",
    ],
    Inches(0.55),
    Inches(3.1),
    Inches(12.0),
    Inches(1.7),
)
add_bullets(
    slide,
    "제어 요청 / 응답 예시",
    [
        "요청: [7E][01][03][22][11][0002][1002 00][CRC16][7E]",
        "응답: [7E][01][04][22][11][0005][00][0005 01 01 00][CRC16][7E]",
        "의미: Relay OFF 요청 / 성공 코드 00 / 실제 Relay_State=OFF",
    ],
    Inches(0.55),
    Inches(4.95),
    Inches(12.0),
    Inches(1.45),
)
add_note(slide, "Control Response에는 성공 여부와 실제 Relay_State를 함께 포함")


# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "예시 해석", "실제 값 파싱 / 3개 단위 분할 예시")
add_table(
    slide,
    Inches(0.45),
    Inches(1.7),
    Inches(12.1),
    Inches(2.35),
    ["예시 데이터", "구성", "해석"],
    [
        ["0002 03 04 00000898", "ITEM_ID=0002 / TYPE=03 / LEN=04 / VALUE=00000898", "Voltage = 2200 -> 220.0V"],
        ["0003 03 04 00001388", "ITEM_ID=0003 / TYPE=03 / LEN=04 / VALUE=00001388", "Current = 5000 -> 50.00A"],
        ["0004 03 04 000186A0", "ITEM_ID=0004 / TYPE=03 / LEN=04 / VALUE=000186A0", "Energy_Acc = 100000Wh"],
        ["0005 01 01 01", "ITEM_ID=0005 / TYPE=01 / LEN=01 / VALUE=01", "Relay_State = ON"],
        ["0006 05 02 0003", "ITEM_ID=0006 / TYPE=05 / LEN=02 / VALUE=0003", "과전류(bit0) + 과전압(bit1)"],
    ],
    font_size=Pt(10),
)
add_bullets(
    slide,
    "4개 항목 응답 예시",
    [
        "1프레임: Voltage / Current / Energy_Acc",
        "2프레임: Relay_State / data_empty / data_empty",
        "응답을 3개 단위로 끊고, 남는 자리는 data_empty(0x0000) 사용",
    ],
    Inches(0.55),
    Inches(4.35),
    Inches(6.1),
    Inches(1.15),
)
add_bullets(
    slide,
    "알람 발생 예시",
    [
        "Alarm_Code = 0x0003이면 bit0, bit1 활성",
        "의미: 과전류 + 과전압 동시 발생",
        "2차 개발 시 Event Report로 즉시 보고 가능",
    ],
    Inches(6.8),
    Inches(4.35),
    Inches(5.35),
    Inches(1.15),
)
add_note(slide, "전압은 raw/10, 전류는 raw/100, 전력량은 raw 그대로 Wh로 해석")


prs.save(OUT_PATH)
with ZipFile(OUT_PATH) as zf:
    slide1_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "지선차단기 - DCU 통신 프로토콜" in slide1_xml
    assert "프레임 형식 / 필드 정의" in slide1_xml

print(OUT_PATH)
print("VERIFY_OK")
