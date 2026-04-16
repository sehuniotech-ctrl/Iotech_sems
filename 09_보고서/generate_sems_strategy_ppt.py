from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"D:\work\15_지선차단기")
OUT = BASE / "08_참고자료" / "SEMS_시장조사_및_SmartLoad_제안서_v0.1.pptx"
IMG_SEMS = BASE / "08_참고자료" / "SEMS_전체전리.png"
IMG_SMARTLOAD = BASE / "08_참고자료" / "지선차단기_기능정리.png"


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


COLORS = {
    "navy": RGBColor(17, 34, 64),
    "blue": RGBColor(34, 88, 160),
    "sky": RGBColor(221, 235, 247),
    "teal": RGBColor(19, 127, 136),
    "green": RGBColor(47, 117, 73),
    "orange": RGBColor(204, 102, 0),
    "red": RGBColor(170, 40, 40),
    "text": RGBColor(41, 41, 41),
    "muted": RGBColor(96, 96, 96),
    "line": RGBColor(210, 218, 228),
    "white": RGBColor(255, 255, 255),
}


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.65))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "맑은 고딕"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = COLORS["navy"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.25), Inches(0.42), Inches(3.3), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.alignment = PP_ALIGN.RIGHT
        run = sub_p.add_run()
        run.text = subtitle
        run.font.name = "맑은 고딕"
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLORS["muted"]

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.02), Inches(12.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.fill.background()


def add_title_slide(title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.25))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "맑은 고딕"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS["navy"]

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = "맑은 고딕"
    run2.font.size = Pt(18)
    run2.font.color.rgb = COLORS["blue"]

    footer_box = slide.shapes.add_textbox(Inches(0.82), Inches(6.55), Inches(11.6), Inches(0.4))
    fp = footer_box.text_frame.paragraphs[0]
    fr = fp.add_run()
    fr.text = footer
    fr.font.name = "맑은 고딕"
    fr.font.size = Pt(11)
    fr.font.color.rgb = COLORS["muted"]

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(5.55), Inches(2.2), Inches(0.6))
    tag.fill.solid()
    tag.fill.fore_color.rgb = COLORS["sky"]
    tag.line.color.rgb = COLORS["blue"]
    tf2 = tag.text_frame
    tf2.clear()
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Draft v0.1"
    r.font.name = "맑은 고딕"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = COLORS["blue"]

    return slide


def add_bullets_slide(title, bullets, right_note=None, subtitle="2026-04-10"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(7.4), Inches(5.55))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        run = p.add_run()
        run.text = bullet
        run.font.name = "맑은 고딕"
        run.font.size = Pt(20 if idx == 0 else 18)
        run.font.color.rgb = COLORS["text"]

    if right_note:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.55), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["sky"]
        card.line.color.rgb = COLORS["line"]
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.clear()
        for idx, line in enumerate(right_note):
            p = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = line
            run.font.name = "맑은 고딕"
            run.font.size = Pt(17 if idx == 0 else 14)
            run.font.bold = idx == 0
            run.font.color.rgb = COLORS["navy"] if idx == 0 else COLORS["text"]

    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, title, "내부 정리 기준")

    for x, section_title, bullets, color in [
        (0.7, left_title, left_bullets, COLORS["blue"]),
        (6.8, right_title, right_bullets, COLORS["teal"]),
    ]:
        head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.35), Inches(5.7), Inches(0.55))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        hp = head.text_frame.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run()
        hr.text = section_title
        hr.font.name = "맑은 고딕"
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = COLORS["white"]

        box = slide.shapes.add_textbox(Inches(x), Inches(2.05), Inches(5.7), Inches(4.6))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = bullet
            r.font.name = "맑은 고딕"
            r.font.size = Pt(16)
            r.font.color.rgb = COLORS["text"]

    return slide


def add_image_slide(title, image_path, caption, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, title, "구조 설명")

    slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.45), width=Inches(6.4))

    cap = slide.shapes.add_textbox(Inches(0.8), Inches(5.95), Inches(6.1), Inches(0.4))
    cp = cap.text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = caption
    cr.font.name = "맑은 고딕"
    cr.font.size = Pt(11)
    cr.font.color.rgb = COLORS["muted"]

    box = slide.shapes.add_textbox(Inches(7.45), Inches(1.55), Inches(5.0), Inches(4.9))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = bullet
        r.font.name = "맑은 고딕"
        r.font.size = Pt(17)
        r.font.color.rgb = COLORS["text"]

    return slide


def add_table_like_slide(title, headers, rows, col_widths, subtitle="시장 비교"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, title, subtitle)

    left = Inches(0.65)
    top = Inches(1.45)
    width = Inches(sum(col_widths))
    row_h = Inches(0.72)

    x = left
    for idx, head in enumerate(headers):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, Inches(col_widths[idx]), row_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["navy"]
        shape.line.color.rgb = COLORS["white"]
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        r.font.name = "맑은 고딕"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]
        x += Inches(col_widths[idx])

    for r_idx, row in enumerate(rows):
        x = left
        for c_idx, cell in enumerate(row):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x,
                top + row_h * (r_idx + 1),
                Inches(col_widths[c_idx]),
                row_h * 1.18,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLORS["sky"] if r_idx % 2 == 0 else COLORS["white"]
            shape.line.color.rgb = COLORS["line"]
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = cell
            r.font.name = "맑은 고딕"
            r.font.size = Pt(12.5)
            r.font.color.rgb = COLORS["text"]
            x += Inches(col_widths[c_idx])

    foot = slide.shapes.add_textbox(Inches(0.75), Inches(6.55), Inches(11.6), Inches(0.35))
    p = foot.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "출처: 한국에너지공단, Schneider Electric, Siemens, Eaton, Leviton 공식 자료 기반 요약"
    r.font.name = "맑은 고딕"
    r.font.size = Pt(10)
    r.font.color.rgb = COLORS["muted"]

    return slide


def add_wbs_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["white"])
    add_header(slide, "WBS 및 마일스톤", "2026-04-10 기준")

    items = [
        ("4/10-4/16", "SEMS 정의 및 범위 확정", COLORS["blue"]),
        ("4/10-4/18", "시장조사 및 타겟 정리", COLORS["teal"]),
        ("4/11-4/22", "기능 및 인터페이스 확정", COLORS["green"]),
        ("4/15-4/29", "HW 회로/핀/전원 구조 확정", COLORS["orange"]),
        ("4/18-5/13", "FW 드라이버/보호/복구 구현", COLORS["blue"]),
        ("5/06-5/20", "데모 시제품 구성", COLORS["teal"]),
        ("5/16-6/05", "시험 및 보정", COLORS["green"]),
        ("5/20-6/10", "제안서/영업 패키지 완성", COLORS["orange"]),
        ("6/01-6/30", "현장 PoC 준비", COLORS["red"]),
    ]

    for idx, (period, text, color) in enumerate(items):
        y = 1.45 + idx * 0.55
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(y), Inches(1.55), Inches(0.38))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        p = tag.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = period
        r.font.name = "맑은 고딕"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]

        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.65), Inches(y), Inches(7.8), Inches(0.38))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["sky"]
        bar.line.color.rgb = COLORS["line"]
        p2 = bar.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = text
        r2.font.name = "맑은 고딕"
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLORS["text"]

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.75), Inches(1.55), Inches(1.9), Inches(4.65))
    note.fill.solid()
    note.fill.fore_color.rgb = COLORS["sky"]
    note.line.color.rgb = COLORS["line"]
    tf = note.text_frame
    tf.word_wrap = True
    lines = [
        "핵심 기준",
        "6월 목표는 양산 완료가 아니라",
        "데모 가능한 시제품 + 제안 패키지 확보로 재정의",
    ]
        # keep concise in card
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = line
        r.font.name = "맑은 고딕"
        r.font.size = Pt(16 if idx == 0 else 13)
        r.font.bold = idx == 0
        r.font.color.rgb = COLORS["navy"] if idx == 0 else COLORS["text"]

    return slide


def add_closing_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["navy"])

    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.1), Inches(11.3), Inches(1.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "결론 및 제안"
    r.font.name = "맑은 고딕"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    bullets = [
        "Smart Load는 단순 차단기가 아니라 SEMS의 말단 실행형 데이터 장치로 포지셔닝해야 한다.",
        "1차는 단상·직접계측·자체판단·차단·DCU연동에 집중하고, 고급 보호 기능은 2차 확장으로 분리한다.",
        "시장 접근은 제품 단품 판매보다 현장 맞춤형 SEMS 구축 패키지 제안이 유리하다.",
    ]

    body = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(10.8), Inches(3.0))
    btf = body.text_frame
    btf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(14)
        r = p.add_run()
        r.text = bullet
        r.font.name = "맑은 고딕"
        r.font.size = Pt(19)
        r.font.color.rgb = COLORS["white"]

    foot = slide.shapes.add_textbox(Inches(1.0), Inches(6.35), Inches(10.5), Inches(0.4))
    p = foot.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Next step: 판매 타겟 확정 / WESYNC 데이터 항목 확정 / 1차 시제품 범위 승인"
    r.font.name = "맑은 고딕"
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(220, 230, 245)

    return slide


add_title_slide(
    "SEMS 시장조사 및 Smart Load 제안",
    "Smart Load(기존 지선차단기) 기반 사업/기술/일정 통합 제안",
    "기준일 2026-04-10 | 내부 공유용 Draft",
)

add_bullets_slide(
    "Executive Summary",
    [
        "SEMS는 단순 계측기가 아니라 계측·보호·분석·절감 운영이 결합된 시스템으로 봐야 합니다.",
        "Smart Load는 SEMS 말단에서 데이터 생성과 현장 차단을 동시에 수행하는 실행형 현장단말입니다.",
        "1차 범위는 단상·직접계측·과전류/과전압·RS-485/DCU 연동에 집중하고, 고급 보호 기능은 2차 확장으로 분리합니다.",
    ],
    [
        "핵심 메시지",
        "제품 단품보다",
        "SEMS 구축용 현장단말",
        "+ 연동 패키지로 제안",
    ],
)

add_image_slide(
    "SEMS 구조에서 Smart Load 위치",
    IMG_SEMS,
    "참고 이미지: 기존 내부 자료 `SEMS_전체전리.png`",
    [
        "전체 구조는 MQTT 서버 - 3상 점검장치 - PLC - DCU - Smart Load 방향으로 정리되어 있습니다.",
        "각 Smart Load는 자기 데이터를 직접 계측하고, 자체 판단으로 차단을 수행한 뒤 DCU에 전달합니다.",
        "즉 Smart Load는 SEMS의 말단 센서가 아니라 실행 기능이 있는 현장단말입니다.",
    ],
)

add_image_slide(
    "Smart Load 제품 개념",
    IMG_SMARTLOAD,
    "참고 이미지: 기존 내부 자료 `지선차단기_기능정리.png`",
    [
        "핵심 기능은 분기회로 계측, 자체 판단, 차단, 상태/이벤트 전송입니다.",
        "대외 명칭은 지선차단기보다 Smart Load를 우선 사용하는 것이 맞습니다.",
        "제품 설명은 단품 차단기보다 현장 데이터 인프라의 한 축으로 잡는 것이 유리합니다.",
    ],
)

add_table_like_slide(
    "시장 및 경쟁 제품 요약",
    ["구분", "대표 제품/근거", "시장 시사점"],
    [
        ["국내 수요", "한국에너지공단 EMS 보급 지원사업", "계측·제어 인프라와 모니터링 시스템을 함께 제안해야 실제 사업화가 쉬움"],
        ["Schneider", "PowerTag", "분기회로별 가시성과 최종 부하 단위 데이터 확보가 중요"],
        ["Siemens", "Sentron Powermind", "측정에서 상태감시·예지보전으로 확장 중"],
        ["Eaton", "Power Xpert BCM", "과부하 경보, 다운타임 감소, 운영 안정성까지 함께 판매"],
        ["Leviton", "VerifEye 7000", "기존 분전반에 붙는 레트로핏과 다양한 프로토콜 지원이 경쟁력"],
    ],
    [2.0, 3.2, 6.8],
)

add_bullets_slide(
    "우리 제품의 특화 포인트",
    [
        "측정 + 자체 판단 + 차단을 하나의 장치에서 수행합니다.",
        "기존 MQTT/PLC/DCU 구조에 붙기 쉬운 현장 맞춤형 단말입니다.",
        "복구정책, 전송주기, 이벤트 정책, 차단 임계치 등 현장별 커스터마이징이 쉽습니다.",
        "1차/2차/3차 확장 로드맵이 분명해 과도한 약속 없이 단계적으로 진입할 수 있습니다.",
    ],
    [
        "추천 포지션",
        "국내 현장 맞춤형",
        "Smart Load 단말",
        "+ SEMS 연동 패키지",
    ],
)

add_two_column_slide(
    "개발 범위 제안",
    "1차 범위",
    [
        "단상 Smart Load",
        "전압, 전류, 주파수, 유효전력, 전력량",
        "과전류/과전압 차단",
        "RS-485 기반 DCU 연동",
        "상태/이벤트/누적값 전송",
    ],
    "2차 확장",
    [
        "피상전력/역률/무효전력 고도화",
        "이력 기반 경보와 패턴 추적",
        "복구 정책 세분화",
        "ZCT 기반 누설전류 확장",
        "고조파/아크/상태진단 검토",
    ],
)

add_bullets_slide(
    "제안 전략",
    [
        "영업 포지션은 차단기 1개 판매보다 현장 데이터 인프라 구축과 절감 운영 기반 제공에 맞춰야 합니다.",
        "1차 판매 타겟은 기존 분전반 데이터 가시성과 보호 기능이 동시에 필요한 중소 현장에 적합합니다.",
        "6월 목표는 양산 완료가 아니라 데모 가능한 시제품과 대외 제안 패키지 확보로 재정의하는 것이 현실적입니다.",
    ],
    [
        "대외 메시지",
        "Smart Load는",
        "SEMS 말단 실행형",
        "현장단말",
    ],
)

add_wbs_slide()

add_table_like_slide(
    "업무분장 및 운영 기준",
    ["영역", "주요 역할", "현재 읽히는 담당"],
    [
        ["PM", "범위 확정, 일정관리, 이슈관리, 대외 조율", "김덕회 팀장"],
        ["상위 서버/플랫폼", "플랫폼 구조, 데이터 활용 방향", "조중협 이사"],
        ["3상 점검장치/연계", "연계 구조, 게이트웨이 영역", "WESYNC"],
        ["Smart Load", "계측, 차단, 통신, 현장단말 구현", "임세훈 주임"],
        ["HW/FW 공통", "회로, 드라이버, 보호/복구, 시험", "세부 재확인 필요"],
    ],
    [1.8, 7.2, 3.0],
    subtitle="업무분장 초안",
)

add_bullets_slide(
    "문서 버전관리 기준",
    [
        "공유 기준 문서는 텍스트 원본과 발표용 PPT를 분리해 관리합니다.",
        "텍스트 기준 문서에서 내용 확정 후 PPT는 파생 산출물로 갱신합니다.",
        "버전은 v0.x 초안, v1.x 내부 확정, v2.x 대외 활용본으로 나눠 운영합니다.",
    ],
    [
        "이번 산출물",
        "텍스트 기준 문서",
        "PPT 제안서 초안",
        "로컬 Git 커밋 완료",
    ],
)

add_closing_slide()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(OUT)
